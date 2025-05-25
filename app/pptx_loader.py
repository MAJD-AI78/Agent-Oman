from pptx import Presentation

def load_pptx_text(path):
    prs = Presentation(path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text
